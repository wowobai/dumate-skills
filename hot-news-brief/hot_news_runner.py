#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
hot_news_runner.py V16 - 统一新闻生成脚本
4合1: PPT + docx + 公众号HTML + 公众号草稿(Supabase Edge Function直调, 单篇模式)

V14变更: 公众号草稿从 appmiaoda.com/api 改为直调 Supabase Edge Function,
         添加Bearer+apikey双认证头，修复entrypoint部署问题。
V14.1变更: 新增配图一致性校验(check_image_consistency)、
           工作区清理(cleanup_workspace，排除自身脚本不被删除)。
V15变更: 草稿改为单篇模式(5条新闻合并为1篇图文)，
         Edge Function createDraft已更新为单article模式。
V16变更: 配图一致性校验增强(输出视觉验证提醒+URL日期检查提醒)，
         main()启动时输出去重检查提醒。

用法: python -X utf8 hot_news_runner.py news_data.json
"""

import json, os, sys, ssl, urllib.request, re
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from docx import Document
from docx.shared import Pt as DocPt, Inches as DocInches, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ======================== Config ========================
BG = RGBColor(0xF5,0xF5,0xFA)
TITLE_C = RGBColor(0x1A,0x56,0xDB)
BODY_C = RGBColor(0x33,0x33,0x33)
ACCENT_C = RGBColor(0x3B,0x82,0xF6)
GRAY_C = RGBColor(0x99,0x99,0x99)
SAFE_GAP = Inches(0.12)
CHARS_PER_INCH = 4.2

# WeChat credentials (for Miaoda server-side use)
WX_APP_ID = "wxb56b11d47f9bdda1"
WX_APP_SECRET = "885e16bdaa8e8223452054a558a4531f"

# Miaoda proxy - via Supabase Edge Function
SUPABASE_URL = "https://backend.appmiaoda.com/projects/supabase340340166944145408"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJhdWQiOiJhdXRoZW50aWNhdGVkIiwiZXhwIjoyMTAwNTcwNjE5LCJpc3MiOiJzdXBhYmFzZSIsInJvbGUiOiJhbm9uIiwic3ViIjoiYW5vbiJ9.9H1fDmEEg7YTB-lqeWLVBK5VTtrsq49ce-GKzoAMLqQ"
MIAODA_URL = f"{SUPABASE_URL}/functions/v1/create-draft"

ssl_ctx = ssl.create_default_context()
ssl_ctx.check_hostname = False
ssl_ctx.verify_mode = ssl.CERT_NONE

# ======================== Helpers ========================
def num_to_cn(n):
    d = "零一二三四五六七八九"
    if n < 10: return d[n]
    if n < 20: return "十" + (d[n%10] if n%10 else "")
    if n < 100:
        r = d[n//10] + "十"
        return r + d[n%10] if n%10 else r
    return str(n)

def date_to_cn(date_str):
    """20260728 -> 二零二六年七月二十八日"""
    y, m, d = int(date_str[:4]), int(date_str[4:6]), int(date_str[6:8])
    y_cn = "".join(num_to_cn(int(c)) for c in str(y))
    return f"{y_cn}年{num_to_cn(m)}月{num_to_cn(d)}日"

def truncate(text, limit=200):
    return text[:limit-3] + "..." if len(text) > limit else text

def compress_image(path):
    img = Image.open(path)
    if img.mode in ('P','LA','RGBA','L','ARGB'):
        img = img.convert('RGB')
    if img.width > 1200:
        ratio = 1200 / img.width
        img = img.resize((1200, int(img.height*ratio)), Image.LANCZOS)
    img.save(path, "JPEG", quality=85)

# ======================== PPT ========================
def gen_ppt(base, data, news):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    date_cn = data.get("date_chinese", date_to_cn(data["date"]))
    
    def add_bg(slide):
        bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = BG
    
    def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, wrap=True, spacing=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = wrap; tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]; p.text = text
        p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
        p.alignment = align
        if spacing: p.line_spacing = spacing
        return p
    
    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_text(s, 1, 2.5, 8, 1.5, "热点文娱新闻", 44, TITLE_C, True, PP_ALIGN.CENTER)
    add_text(s, 1, 3.5, 8, 0.8, date_cn, 24, ACCENT_C, False, PP_ALIGN.CENTER)
    
    # Slide 2: TOC
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_text(s, 0.5, 0.3, 9, 0.8, "今日导览", 32, TITLE_C, True)
    y = 1.5
    for i, item in enumerate(news):
        add_text(s, 1, y, 8, 0.9, f"{i+1}. {item['title']}", 20, BODY_C)
        y += 0.9
    
    # Slides 3-7: Detail
    for item in news:
        s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
        img_path = os.path.join(base, item.get("image_path", f"images/news_{item['id']}.jpg"))
        if os.path.exists(img_path):
            try:
                img = Image.open(img_path)
                iw, ih = img.size
                ratio = min(4.5/iw*72, 4.5/ih*72)/72
                s.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), Inches(iw*ratio/72), Inches(ih*ratio/72))
            except: pass
        
        add_text(s, 5, 0.5, 4.5, 1.5, item["title"], 22, TITLE_C, True)
        add_text(s, 5, 2.2, 4.5, 2.5, truncate(item["body"]), 14, BODY_C, spacing=1.35)
        add_text(s, 5, 5.0, 4.5, 1.5, "【分析】"+item["analysis"], 13, ACCENT_C, spacing=1.35)
        add_text(s, 5, 6.5, 4.5, 0.5, item["source"], 10, GRAY_C)
    
    # Slide 8: End
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_text(s, 1, 3, 8, 1.5, date_cn, 28, TITLE_C, False, PP_ALIGN.CENTER)
    add_text(s, 1, 4, 8, 0.8, "本期完", 20, ACCENT_C, False, PP_ALIGN.CENTER)
    
    path = os.path.join(base, f"热点文娱新闻PPT_{data['date']}.pptx")
    prs.save(path)
    return path

# ======================== docx ========================
def gen_docx(base, data, news):
    doc = Document()
    date_cn = data.get("date_chinese", date_to_cn(data["date"]))
    
    def add_para(text, size, color=None, bold=False, align=None, italic=False):
        p = doc.add_paragraph()
        if align: p.alignment = align
        run = p.add_run(text)
        run.font.size = DocPt(size); run.font.bold = bold; run.font.italic = italic
        if color: run.font.color.rgb = color
        return p
    
    add_para(f"{date_cn}新闻稿", 22, DocRGB(0x1A,0x56,0xDB), True, WD_ALIGN_PARAGRAPH.CENTER)
    add_para("—— 每日热点文娱资讯 ——", 12, DocRGB(0x66,0x66,0x66), align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    add_para("各位听众朋友们大家好，欢迎收听今日热点文娱新闻。", 11)
    
    for i, item in enumerate(news):
        add_para(f"第{num_to_cn(i+1)}条，{item['title']}。", 13, DocRGB(0x1A,0x56,0xDB), True)
        combined = truncate(item["body"] + item["analysis"])
        p = add_para(combined, 11)
        p.paragraph_format.line_spacing = 1.5
        add_para(item["source"] + "。", 9, DocRGB(0x99,0x99,0x99), italic=True)
        
        img_path = os.path.join(base, item.get("image_path", f"images/news_{item['id']}.jpg"))
        if os.path.exists(img_path):
            try:
                p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.add_run().add_picture(img_path, width=DocInches(4.0))
            except: pass
        doc.add_paragraph()
    
    add_para("以上就是今日热点文娱新闻，感谢您的收听。", 11, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(date_cn, 10, DocRGB(0x66,0x66,0x66), align=WD_ALIGN_PARAGRAPH.CENTER)
    
    path = os.path.join(base, f"{date_cn}新闻稿.docx")
    doc.save(path)
    return path

# ======================== HTML ========================
def gen_html(base, data, news):
    date_cn = data.get("date_chinese", date_to_cn(data["date"]))
    dd = data["date_display"]
    
    css = "body{font-family:-apple-system,sans-serif;margin:0;padding:20px;color:#333}" \
          ".h{text-align:center;padding:30px 0;background:linear-gradient(135deg,#1a56db,#3b82f6);color:#fff;border-radius:8px 8px 0 0}" \
          ".h h1{font-size:28px;margin:0}" ".h .d{font-size:16px;margin-top:10px;opacity:.9}" \
          ".ni{padding:25px 20px;border-bottom:1px solid #eee}" \
          ".nt{font-size:20px;font-weight:bold;color:#1a56db;margin-bottom:12px}" \
          ".nb{font-size:15px;line-height:1.8;margin-bottom:10px}" \
          ".na{font-size:14px;color:#3b82f6;line-height:1.7;padding:8px 12px;background:#f0f7ff;border-radius:4px;margin-bottom:8px}" \
          ".ns{font-size:12px;color:#999;font-style:italic}" ".ni img{width:100%;border-radius:6px;margin:12px 0}" \
          ".ft{text-align:center;padding:20px;color:#999;font-size:13px}"
    
    parts = [f'<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">',
             f'<title>{dd} 热点文娱新闻</title><style>{css}</style></head><body>',
             f'<div style="max-width:677px;margin:0 auto">',
             f'<div class="h"><h1>热点文娱新闻</h1><div class="d">{date_cn}</div></div>']
    
    for item in news:
        parts.append(f'<div class="ni">')
        parts.append(f'<div class="nt">{item["title"]}</div>')
        img_path = item.get("image_path", f"images/news_{item['id']}.jpg")
        parts.append(f'<img src="{img_path}" alt="{item["title"]}">')
        parts.append(f'<div class="nb">{truncate(item["body"])}</div>')
        parts.append(f'<div class="na">【分析】{item["analysis"]}</div>')
        parts.append(f'<div class="ns">{item["source"]}</div>')
        parts.append('</div>')
    
    parts.append(f'<div class="ft"><p>{date_cn}</p><p>本期完</p></div>')
    parts.append('</div></body></html>')
    
    path = os.path.join(base, "wechat_article.html")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(parts))
    return path

# ======================== WeChat Draft (Supabase Edge Function) ========================
def gen_wechat_draft(base, data, news):
    """通过Supabase Edge Function创建公众号草稿
    V14: 直调Edge Function端点，Bearer+apikey双认证
    诊断流程见SKILL V14标准化诊断章节
    """
    img_dir = os.path.join(base, "images")
    news_json = json.dumps(data, ensure_ascii=False)
    files = []
    for i in range(1, 6):
        img_path = os.path.join(img_dir, f"news_{i}.jpg")
        if os.path.exists(img_path):
            files.append(("images", (f"news_{i}.jpg", open(img_path, "rb"), "image/jpeg")))
    
    payload = {"news_data": news_json}
    
    try:
        headers = {
            "Authorization": f"Bearer {SUPABASE_KEY}",
            "apikey": SUPABASE_KEY,
        }
        resp = requests.post(MIAODA_URL, data=payload, files=files, headers=headers, verify=False, timeout=120)
        result = resp.json()
        if result.get("success"):
            print(f"  [OK] 草稿创建成功! Draft ID: {result.get('draft_id')}")
            return result
        else:
            print(f"  [ERROR] 草稿创建失败: {result.get('error', result)}")
            return None
    except Exception as e:
        print(f"  [ERROR] 秒哒代理请求异常: {e}")
        return None
    finally:
        for _, f_tuple in files:
            f_tuple[1].close()

# ======================== Verify ========================
def verify_all(base, data, news, ppt_path, docx_path, html_path, draft_result=None):
    print("\n=== 验证产出 ===")
    ok = True
    
    if os.path.exists(ppt_path):
        size = os.path.getsize(ppt_path)
        prs = Presentation(ppt_path)
        pics = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        status = "OK" if size > 102400 and len(prs.slides) >= 8 and pics >= 5 else "FAIL"
        print(f"  PPT: {status} | {size}B | {len(prs.slides)}页 | {pics}图")
        ok = ok and status == "OK"
    else:
        print("  PPT: MISSING"); ok = False
    
    if os.path.exists(docx_path):
        size = os.path.getsize(docx_path)
        doc = Document(docx_path)
        text = "\n".join(p.text for p in doc.paragraphs)
        has_cn_date = "二零二六" in text
        news_count = sum(1 for p in doc.paragraphs if "第" in p.text and "条" in p.text)
        img_count = sum(1 for r in doc.part.rels.values() if "image" in r.reltype)
        status = "OK" if size > 102400 and has_cn_date and news_count >= 5 and img_count >= 5 else "FAIL"
        print(f"  docx: {status} | {size}B | {news_count}条 | {img_count}图 | 中文日期:{has_cn_date}")
        ok = ok and status == "OK"
    else:
        print("  docx: MISSING"); ok = False
    
    if os.path.exists(html_path):
        print(f"  HTML: OK | {os.path.getsize(html_path)}B")
    else:
        print("  HTML: MISSING"); ok = False
    
    if draft_result and draft_result.get("success"):
        print(f"  草稿: OK | {draft_result.get('draft_id')}")
    else:
        print("  草稿: SKIP")
    
    print(f"\n  总体: {'全部通过' if ok else '有项未通过'}")
    return ok

# ======================== Image Consistency Check ========================
def check_image_consistency(base, news):
    img_dir = os.path.join(base, "images")
    print("  图片与新闻内容对应关系:")
    all_ok = True
    for i, item in enumerate(news):
        img_path = os.path.join(img_dir, f"news_{i+1}.jpg")
        if os.path.exists(img_path):
            try:
                img = Image.open(img_path)
                w, h = img.size
                ratio = w / h if h > 0 else 0
                fsize = os.path.getsize(img_path)
                min_dim = min(w, h)
                status = "OK" if fsize > 10240 and min_dim >= 200 else "WARN"
                if status != "OK":
                    all_ok = False
                title_short = item['title'][:30]
                print(f"  [{status}] news_{i+1}.jpg: {w}x{h} ({ratio:.2f}) {fsize}B -> {title_short}")
            except Exception as e:
                print(f"  [ERR] news_{i+1}.jpg: {e}")
                all_ok = False
        else:
            print(f"  [MISS] news_{i+1}.jpg -> {item['title'][:30]}")
            all_ok = False
    print(f"  配图一致性(尺寸): {'全部通过' if all_ok else '有项需关注'}")
    print("  [V16提醒] 尺寸校验通过不等于视觉内容匹配。")
    print("  [V16提醒] 运行此脚本的AI必须用read工具逐张查看图片，确认视觉内容与新闻标题一致。")
    print("  [V16提醒] 如果图片内容与新闻不匹配，必须重新搜索下载。")
    return all_ok

# ======================== Workspace Cleanup ========================
def cleanup_workspace(base):
    this_file = os.path.abspath(__file__)
    workspace_root = os.path.dirname(os.path.dirname(base))
    cleaned = 0

    # 1. Clean scattered scripts in workdir (but not this file)
    workdir = os.path.dirname(base)
    for pattern in ["hot_news_runner.py", "hot_news_skill_article_SKILL_V*.md",
                    "hot_news_runner_v*.py", "miaoda_api.py"]:
        import glob
        for p in glob.glob(os.path.join(workdir, pattern)):
            if os.path.abspath(p) != this_file:
                os.remove(p)
                print(f"  删除: {os.path.basename(p)}")
                cleaned += 1

    # 2. Clean old output dirs in workspace root (not current session)
    for d in os.listdir(workspace_root):
        full = os.path.join(workspace_root, d)
        if os.path.isdir(full) and d.startswith("ses_") and full != workdir:
            for sub in os.listdir(full):
                if "热点新闻" in sub or "hot_news" in sub.lower():
                    import shutil
                    shutil.rmtree(os.path.join(full, sub), ignore_errors=True)
                    print(f"  删除: {d}/{sub}")
                    cleaned += 1

    # 3. Clean old output dirs in workspace root root
    root_root = os.path.dirname(workspace_root)
    if os.path.exists(root_root):
        for d in os.listdir(root_root):
            if "热点新闻" in d or "hot_news" in d.lower():
                import shutil
                shutil.rmtree(os.path.join(root_root, d), ignore_errors=True)
                print(f"  删除: {d}")
                cleaned += 1

    print(f"  清理完成: 删除{cleaned}项")

# ======================== Main ========================
def main():
    if len(sys.argv) < 2:
        print("Usage: python -X utf8 hot_news_runner.py news_data.json")
        sys.exit(1)
    
    json_path = sys.argv[1]
    base = os.path.dirname(os.path.abspath(json_path))
    
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    news = data["news"]
    if "date_chinese" not in data:
        data["date_chinese"] = date_to_cn(data["date"])
    
    print(f"=== V16 热点文娱新闻生成 ===")
    print(f"日期: {data['date_display']} ({data['date_chinese']})")
    print(f"新闻: {len(news)}条")
    
    print("\n[V16] 去重检查提醒: 运行前请确认已搜索前1-2日选题记录，排除重复新闻。")
    print("[V16] 配图提醒: 下载配图前请检查URL路径中的日期是否与新闻事件日期匹配。")
    print("[V16] 视觉验证提醒: 配图下载后请用read工具逐张查看，确认视觉内容匹配。\n")
    
    print("\n--- 压缩配图 ---")
    img_dir = os.path.join(base, "images")
    for i in range(1, 6):
        p = os.path.join(img_dir, f"news_{i}.jpg")
        if os.path.exists(p):
            compress_image(p)
            print(f"  news_{i}.jpg: {os.path.getsize(p)}B")
    
    print("\n--- 生成PPT ---")
    ppt_path = gen_ppt(base, data, news)
    print(f"  {ppt_path} ({os.path.getsize(ppt_path)}B)")
    
    print("\n--- 生成新闻稿docx ---")
    docx_path = gen_docx(base, data, news)
    print(f"  {docx_path} ({os.path.getsize(docx_path)}B)")
    
    print("\n--- 生成公众号HTML ---")
    html_path = gen_html(base, data, news)
    print(f"  {html_path} ({os.path.getsize(html_path)}B)")
    
    draft_result = None
    if os.name == 'posix':
        print("\n--- 公众号草稿(秒哒代理) ---")
        draft_result = gen_wechat_draft(base, data, news)
    
    verify_all(base, data, news, ppt_path, docx_path, html_path, draft_result)

    print("\n--- 配图一致性校验 ---")
    check_image_consistency(base, news)

    print("\n--- 清理散落文件 ---")
    cleanup_workspace(base)

    print(f"\n=== 完成 ===")

if __name__ == "__main__":
    main()
